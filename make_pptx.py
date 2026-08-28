# Собрать черновик защитной презентации из готовых картинок.
# Тексты и картинки по слайдам заданы ниже; каждой картинки, которой нет на диске,
# скрипт просто не вставляет (и пишет заметку). Дальше причёсывается в PowerPoint.
#
#   pip install python-pptx        # если не стоит
#   python make_pptx.py            # -> presentation.pptx
#   python make_pptx.py --img-dir outputs/figure --sthelar-dir outputs/sthelar --out presentation.pptx

import argparse
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

BLUE = RGBColor(0x1F, 0x49, 0x7D)
DARK = RGBColor(0x22, 0x22, 0x22)
GREY = RGBColor(0x66, 0x66, 0x66)
EMU_IN = 914400

# каждый слайд: (заголовок, [буллеты], [картинки], заметки докладчику)
SLIDES = [
    ("Автоматическая сегментация тканевых зон рака яичника по H&E",
     ["Шилина В.В.", "Научный руководитель: Казакова А.Н.",
      "МФТИ · ФНКЦ ФХМ ФМБА России · 2026"],
     [], "Титульный слайд."),

    ("Актуальность",
     ["HGSC — самый частый и агрессивный рак яичника; ~⅓ пациенток резистентны к платине",
      "Нужны простые биомаркёры по обычному стеклу",
      "TSR (доля стромы) считывается с H&E и связан с химиорезистентностью при HGSC",
      "Но оценивают TSR вручную — зависит от эксперта"],
     ["a3_zoom_he.png"],
     "Треть пациенток резистентны к платине. TSR связан с этим, но считается вручную."),

    ("Строма неоднородна: гормональная и матриксная",
     ["Гормонпродуцирующая строма — стероидогенные C7+ CAFs",
      "Матриксная (десмопластическая) — mCAFs, откладывают коллаген",
      "Разный биологический и прогностический смысл",
      "(картинки из ЛИТОБЗОР.pptx, слайды 34–35)"],
     [],
     "Работа лаборатории: строма делится на функциональные типы — гормональную и матриксную."),

    ("Проблема",
     ["Подтипы стромы на H&E глазом неразличимы",
      "Пространственная транскриптомика их видит, но дорога и есть на единицах срезов",
      "Обычных H&E — огромные архивы",
      "Идея: обучить на немногих срезах с транскриптомикой, применять на обычном H&E"],
     ["a3_zoom_he.png"],
     "Молекулярно подтипы различимы, морфологически — нет. Учим на транскриптомике, применяем без неё."),

    ("Цель и задачи",
     ["Цель: инструмент сегментации 5 тканевых зон рака яичника по H&E",
      "1. Собрать парные данные H&E + Xenium",
      "2. Молекулярные метки 5 зон: опухоль, гормональная/матриксная строма, иммунные, сосуды",
      "3. Предобработка WSI: чтение, нарезка, нормализация окраски",
      "4. Признаки замороженным UNI",
      "5. Граф клеток + две графовые сети (слой 1 опухоль/строма, слой 2 подтип)",
      "6. Оценка leave-one-slide-out; сравнение с попатчевым декодером",
      "7. Внешняя валидация; применение к TCGA без транскриптомики"],
     [],
     "Цель — инструмент по H&E. Ключевое — графовая модель по клеткам, учитывающая окружение."),

    ("Метрики",
     ["IoU и Dice — по клеткам (честное сравнение граф vs патч)",
      "macro-F1 и сбалансированная точность — из-за перекоса классов",
      "Для TSR (непрерывная величина): ICC и Бланд–Альтман",
      "Доверительные интервалы — бутстрепом по срезам"],
     [],
     "Единица оценки — клетка. Для доли стромы — метрики согласия, как в клинике."),

    ("Литобзор: фундаментальные модели патологии",
     ["Замороженный кодировщик + лёгкая голова — меньше потребность в разметке",
      "UNI (ViT-L/16, самообучение на 100 000+ слайдов)",
      "На раке яичника UNI ≈ лучшая модель при вчетверо меньших затратах → выбор UNI"],
     ["b1_patch_grid.png"],
     "Основа — UNI, обучен самообучением на сотнях тысяч слайдов. Поверх — лёгкая голова."),

    ("Литобзор: транскриптомика как источник меток",
     ["Xenium — клеточное разрешение, парный H&E",
      "Деконволюция / выравнивание на атласы → метки типов клеток",
      "STHELAR — 11 млн размеченных клеток (для внешней проверки)"],
     ["a2_cells_thumb.png"],
     "Метки берём из транскриптомики, а не у патолога. STHELAR — публичный набор для внешней проверки."),

    ("Литобзор: клеточные графы и ближайший аналог",
     ["Графы клеток: CGC-Net, HACT-Net, SlideGraph+ (узел = клетка, граф = ткань)",
      "SpatialFusion — замороженные кодировщики + граф, апробирован на раке яичника",
      "Пробел: сегментация компартментов графом на молекулярных метках ранее не описана"],
     [],
     "Клеточные графы — самостоятельная линия. Сегментации компартментов с делением стромы до нас не делали."),

    ("Литобзор: детекция клеток и клиника",
     ["StarDist, HoVer-Net — детекция ядер (StarDist — детектор при применении)",
      "STAR — автоматический TSR по H&E улучшает прогноз выживаемости"],
     ["b2_graph_zoom.png"],
     "Клетки находит StarDist. STAR показывает, что автоматический TSR клинически полезен."),

    ("Данные",
     ["3 среза яичника Xenium + парный H&E: ovary_prime, ovary2, ovary3",
      "Качество совмещения Xenium↔H&E: Dice 0.99 / 0.99 / 0.82",
      "ovary3 — деформация, в обучении, но не в валидации"],
     ["a1_he_thumb.png"],
     "Три среза Xenium с H&E. Совмещение координат критично: третий деформирован, в валидацию не беру."),

    ("Метод — общая схема",
     ["Слева — обучение на транскриптомике (один раз), справа — применение на H&E"],
     ["figure_pipeline.png"],
     "Весь конвейер: слева обучение, справа применение. Разберу по частям."),

    ("Построение разметки (панель a)",
     ["H&E → типы клеток из Xenium → совмещение координат → сплошная маска-истина"],
     ["a1_he_thumb.png", "a2_cells_thumb.png", "a4_zoom_cells.png", "a5_voronoi.png"],
     "Метки без патолога: типы клеток из Xenium, совмещение с H&E, маска по ближайшей клетке."),

    ("Признаки и граф клеток (панель b)",
     ["Патчи → UNI (заморожен) → клетка берёт свой токен → граф (ребро ≤ 50 мкм, k=8)"],
     ["b1_patch_grid.png", "b2_graph_zoom.png"],
     "Кодируем UNI один раз на патч; клетка берёт свой токен. Строим граф соседей."),

    ("Две графовые сети (панель c)",
     ["Слой 1: опухоль/строма (все клетки)", "Слой 2: подтип стромы (только строма)",
      "GraphConv ×2 → контекст ≈ 100 мкм"],
     ["c_layer1_whole.png", "c3_layer2_whole.png"],
     "Графовая свёртка = смесь себя и соседей; два слоя дают ~100 мкм контекста."),

    ("Результат: слой 1 (опухоль/строма)",
     ["Кросс-валидация по срезам: граф mIoU 0.836 против U-Net 0.750",
      "Dice: опухоль 0.92, строма 0.90",
      "Те же признаки UNI → выигрыш от учёта соседей, а не размера сети"],
     ["c_layer1_whole.png"],
     "Главный результат: 0.836 против 0.750 при тех же признаках — чистый вклад контекста."),

    ("Результат: слой 2 (подтипы стромы)",
     ["Внутри среза: матриксная/иммунная/сосудистая Dice 0.65–0.71, гормональная 0.31",
      "Между пациентами: компартменты CAF/иммунный/сосудистый ~0.55",
      "Гормональная/матриксная — уверенно внутри среза; между пациентами — граница на 3 срезах"],
     ["c3_layer2_whole.png"],
     "Три компартмента переносятся между пациентами. Тонкое деление CAF — внутри среза."),

    ("Применение к TCGA без транскриптомики (панель d)",
     ["StarDist → граф → карта пяти зон → TSR и состав стромы автоматически",
      "Пример применения; точность на TCGA не размечена"],
     ["d1_tcga_he.png", "d2_zonemap.png"],
     "На TCGA транскриптомики нет: клетки находит StarDist, дальше тот же граф. Выход — карта зон и TSR."),

    ("Внешняя валидация: STHELAR",
     ["Независимый датасет, граф на нём не обучался",
      "Компартменты переносятся: иммунный Dice ~0.76, сосудистый ~0.81",
      "Опухоль/строма на мелких изолированных тайлах хуже — контекст ломается"],
     ["ovary_s0_0.png", "ovary_s1_0.png"],
     "На независимом STHELAR компартменты стромы переносятся без дообучения."),

    ("Было → стало",
     ["Было: попатчевый декодер, потолок ~0.71–0.75; отвергнуты Visium, суперпиксели, аугментация",
      "Стало: граф клеток 0.836; деление стромы на подтипы"],
     [],
     "Потолок попатчевого подхода пробил переход на граф клеток."),

    ("Что сделано / статус",
     ["✅ Модель: оба слоя, инференс на любом H&E (StarDist→граф→зоны)",
      "✅ Внутренняя валидация (leave-one-slide-out)",
      "✅ Внешняя валидация STHELAR (компартменты)",
      "✅ Схема метода, применение к TCGA",
      "⏳ В работе: UBC-OCEAN, TSR-vs-прогноз на TCGA, обёртка для R"],
     [],
     "Готовы обе сети, полный инференс, внутренняя и внешняя валидация. В работе — прогноз и маски патолога."),

    ("Открытые вопросы и дальнейшее",
     ["Расхождение обучение/инференс по нормализации окраски (Macenko) — проверяем ±Macenko",
      "Достаточность 3 срезов для межпациентного деления гормональная/матриксная",
      "Кривая контекстного радиуса — прямое измерение роли соседства",
      "Вложенный Кокс на TCGA: добавляют ли подтипы прогноз сверх TSR"],
     [],
     "Дальше — закрыть Macenko, измерить кривую контекстного радиуса, проверить прогноз подтипов сверх TSR."),

    ("Выводы / новизна",
     ["1. Впервые — деление стромы яичника на функциональные подтипы по обычному H&E",
      "2. Показано количественно: принадлежность клетки определяется окружением (граф > патч)",
      "3. Перенос компартментов стромы на независимый датасет без дообучения"],
     [],
     "Три вывода: первое деление стромы по H&E; контекст решает; внешний перенос компартментов."),
]


def place_images(slide, paths, left, top, width, height):
    from PIL import Image
    paths = [p for p in paths if p]
    if not paths:
        return
    n = len(paths)
    cols = 1 if n == 1 else (2 if n <= 4 else 3)
    rows = (n + cols - 1) // cols
    cw, ch = width / cols, height / rows
    gap = Inches(0.12)
    for i, p in enumerate(paths):
        r, c = divmod(i, cols)
        iw, ih = Image.open(p).size
        bw, bh = cw - gap, ch - gap
        s = min(bw / iw, bh / ih)
        w, h = iw * s, ih * s
        x = left + c * cw + (cw - w) / 2
        y = top + r * ch + (ch - h) / 2
        slide.shapes.add_picture(str(p), int(x), int(y), int(w), int(h))


def add_title_bar(slide, text, sw):
    bar = slide.shapes.add_shape(1, 0, 0, sw, Inches(0.95))  # 1 = rectangle
    bar.fill.solid(); bar.fill.fore_color.rgb = BLUE
    bar.line.fill.background()
    tf = bar.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.4); tf.margin_top = Inches(0.12)
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(24); p.font.bold = True; p.font.color.rgb = RGBColor(255, 255, 255)


def add_bullets(slide, bullets, left, top, width, height):
    if not bullets:
        return
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame; tf.word_wrap = True
    sz = 14 if len(bullets) > 6 else 17
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ("• " + b) if not b[0].isdigit() and b[0] not in "✅⏳(" else b
        p.font.size = Pt(sz); p.font.color.rgb = DARK
        p.space_after = Pt(8)


def resolve(name, img_dir, sth_dir):
    if not name:
        return None
    for d in (img_dir, sth_dir):
        p = d / name
        if p.exists():
            return p
    return None


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--img-dir", default="outputs/figure")
    ap.add_argument("--sthelar-dir", default="outputs/sthelar")
    ap.add_argument("--out", default="presentation.pptx")
    args = ap.parse_args()

    img_dir, sth_dir = Path(args.img_dir), Path(args.sthelar_dir)
    prs = Presentation()
    prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    sw, sh = prs.slide_width, prs.slide_height
    blank = prs.slide_layouts[6]

    missing = []
    for i, (title, bullets, imgs, notes) in enumerate(SLIDES, 1):
        slide = prs.slides.add_slide(blank)
        found = []
        for nm in imgs:
            p = resolve(nm, img_dir, sth_dir)
            (found.append(p) if p else missing.append((i, nm)))
        if i == 1:  # титул — по центру, без шапки
            tb = slide.shapes.add_textbox(Inches(1), Inches(2.3), Inches(11.3), Inches(3))
            tf = tb.text_frame; tf.word_wrap = True
            p0 = tf.paragraphs[0]; p0.text = title
            p0.font.size = Pt(30); p0.font.bold = True; p0.font.color.rgb = BLUE
            p0.alignment = PP_ALIGN.CENTER
            for b in bullets:
                p = tf.add_paragraph(); p.text = b; p.font.size = Pt(18)
                p.alignment = PP_ALIGN.CENTER; p.font.color.rgb = GREY
        else:
            add_title_bar(slide, title, sw)
            top = Inches(1.2)
            if found and bullets:
                add_bullets(slide, bullets, Inches(0.5), top, Inches(6.0), Inches(5.8))
                place_images(slide, found, Inches(6.7), top, Inches(6.2), Inches(5.9))
            elif found:
                place_images(slide, found, Inches(0.5), top, Inches(12.3), Inches(5.9))
            else:
                add_bullets(slide, bullets, Inches(0.6), top, Inches(12.0), Inches(5.8))
        # заметки докладчику
        slide.notes_slide.notes_text_frame.text = notes

    prs.save(args.out)
    print("сохранено:", args.out, "| слайдов:", len(SLIDES))
    if missing:
        print("не найдены картинки (слайд, файл) — вставишь вручную:")
        for i, nm in missing:
            print(f"  слайд {i}: {nm}")


if __name__ == "__main__":
    main()
